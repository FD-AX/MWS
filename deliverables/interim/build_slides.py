"""Сборка промежуточной презентации (PPTX): гипотезы → эксперименты → цифры → архитектура.
Каждый этап конвейера — отдельный слайд: что делает, зачем, чем доказано, пример находки.

    python deliverables/interim/build_slides.py  → deliverables/interim/TZ_Review_interim.pptx
Графики и схема рисуются matplotlib'ом в PNG рядом (charts/).
"""
from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

HERE = Path(__file__).resolve().parent
OUT = HERE / "TZ_Review_interim.pptx"
CH = HERE / "charts"
CH.mkdir(exist_ok=True)

INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
PAL = ["#9ca3af", "#2563eb", "#059669", "#d97706", "#7c3aed", "#dc2626"]

plt.rcParams.update({"font.size": 11, "axes.spines.top": False, "axes.spines.right": False})


# ---------------------------------------------------------------- графики
def chart_architecture(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(13, 6.2), dpi=150)
    ax.set_xlim(0, 13); ax.set_ylim(0, 6.2); ax.axis("off")

    def box(x, y, w, h, title, sub="", color="#eef2ff", edge="#2563eb"):
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.12", fc=color, ec=edge, lw=1.4))
        ax.text(x + w / 2, y + h - 0.32, title, ha="center", va="center", fontsize=11, weight="bold", color="#1f2328")
        if sub:
            ax.text(x + w / 2, y + h / 2 - 0.22, sub, ha="center", va="center", fontsize=8.6, color="#374151", linespacing=1.35)

    def arrow(x1, y1, x2, y2, label=""):
        ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>", mutation_scale=14, lw=1.2, color="#6b7280"))
        if label:
            ax.text((x1 + x2) / 2, (y1 + y2) / 2 + 0.16, label, ha="center", fontsize=8, color="#6b7280")

    box(0.2, 4.2, 2.1, 1.4, "Аналитик", "UI / Swagger\nmd · txt · docx · pdf", color="#f9fafb", edge="#9ca3af")
    box(2.8, 4.2, 2.1, 1.4, "api", "FastAPI · UI · история\n:18080")
    box(2.8, 2.2, 2.1, 1.4, "docs", "pdf/docx → markdown\nсекции, таблицы")
    box(5.4, 4.2, 2.1, 1.4, "RabbitMQ", "review.jobs (durable)\nDLQ review.dead")
    box(8.0, 4.2, 2.4, 1.4, "worker", "prefetch 1 · ack после записи\nидемпотентность · прогресс")
    box(10.9, 4.2, 2.0, 1.4, "LLM", "gpt-oss-120b (vLLM/ollama)\nGPT-5.5 (референс)", color="#fef3c7", edge="#d97706")
    box(5.4, 2.2, 2.1, 1.4, "Postgres", "documents · reviews\nfindings · feedback", color="#ecfdf5", edge="#059669")
    box(8.0, 2.2, 2.4, 1.4, "Prometheus + Grafana", "поток, p95, вызовы, токены,\nнаходки по классам, UNKNOWN", color="#ecfdf5", edge="#059669")
    box(0.2, 0.15, 12.7, 1.6, "Конвейер tz_review (внутри worker)",
        "детерминированный слой (шаблон МТС, паттерны, граф сущностей) → чеклист 27 слотов (OK / MISSING / UNCLEAR / NA) → "
        "согласованность разделов → взгляд разработчика\n→ энтропия / логит-зонд → верификация цитат → критик "
        "(официальные пункты и детерминированные находки не режутся)", color="#f3f4f6", edge="#4b5563")
    arrow(2.3, 4.9, 2.8, 4.9); arrow(3.85, 4.2, 3.85, 3.6, "normalize"); arrow(4.9, 4.9, 5.4, 4.9, "job")
    arrow(7.5, 4.9, 8.0, 4.9); arrow(10.4, 4.9, 10.9, 4.9, "OpenAI API"); arrow(9.2, 4.2, 6.6, 3.6, "результат")
    arrow(3.85, 4.2, 5.4, 3.3, "история"); arrow(9.2, 4.2, 9.2, 3.6, "/metrics"); arrow(10.4, 4.5, 12.0, 1.75, "review()")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


def chart_pipeline(path: Path) -> None:
    """Конвейер как лента: этап → что даёт (цифра)."""
    stages = [
        ("Детерминированный\nслой", "шаблон МТС, паттерны,\nпустые разделы", "F-класс 6/6\nшум 0"),
        ("Граф\nсущностей", "поля-фантомы, дрейф\nимён, битые ссылки", "+3 synth, +critical doc3\n0 ложных, 0 LLM"),
        ("Чеклист\n27 слотов", "8 пунктов МТС,\nOK/MISSING/UNCLEAR/NA", "пункты МТС 7/7\n(было 3/7)"),
        ("Согласованность\nразделов", "противоречия,\nдрейф терминов", "D-класс 10/11"),
        ("Взгляд\nразработчика", "«что придётся\nугадывать»", "уникальные A3/B9\n(fallback, поля времени)"),
        ("Энтропия /\nлогит-зонд", "неоднозначность\nчислом", "doc3 16/16,\nанафора только здесь"),
        ("Верификация\nцитат", "цитата ⊂ документ,\nиначе — отбрасываем", "anchoring 96–100%\n(без неё 47%)"),
        ("Критик", "ранжирует, дедуп,\nне режет пункты МТС", "порог 4: −40% хвоста,\nrecall держится"),
    ]
    fig, ax = plt.subplots(figsize=(13.2, 3.6), dpi=150)
    ax.set_xlim(0, 13.2); ax.set_ylim(0, 3.6); ax.axis("off")
    w, gap = 1.5, 0.14
    for i, (t, sub, ev) in enumerate(stages):
        x = 0.1 + i * (w + gap)
        col = "#eef2ff" if i not in (0, 1) else "#f3f4f6"
        edge = "#2563eb" if i not in (0, 1) else "#4b5563"
        ax.add_patch(FancyBboxPatch((x, 1.55), w, 1.85, boxstyle="round,pad=0.02,rounding_size=0.1", fc=col, ec=edge, lw=1.3))
        ax.text(x + w / 2, 3.08, t, ha="center", va="center", fontsize=9.6, weight="bold", color="#1f2328")
        ax.text(x + w / 2, 2.25, sub, ha="center", va="center", fontsize=7.8, color="#374151", linespacing=1.3)
        ax.add_patch(FancyBboxPatch((x, 0.25), w, 1.05, boxstyle="round,pad=0.02,rounding_size=0.1", fc="#ecfdf5", ec="#059669", lw=1.1))
        ax.text(x + w / 2, 0.77, ev, ha="center", va="center", fontsize=7.8, color="#065f46", linespacing=1.3)
        if i < len(stages) - 1:
            ax.add_patch(FancyArrowPatch((x + w, 2.45), (x + w + gap, 2.45), arrowstyle="-|>", mutation_scale=10, lw=1, color="#6b7280"))
    ax.text(0.1, 1.42, "что доказано (журнал экспериментов):", fontsize=8, color="#065f46", style="italic")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


def chart_recall(path: Path) -> None:
    variants = ["Детерм. слой", "Промпт GPT-5.5", "Конвейер GPT-5.5", "+ энтропия", "+ логит-зонд", "Конвейер 120b"]
    targets = ["doc3 витрина (16)", "пункты МТС, synth (13)", "synth v1 (12)", "v2hard (8)"]
    data = {  # доли recall
        "doc3 витрина (16)":        [3 / 16, 11 / 16, 15 / 16, 16 / 16, 16 / 16, 13 / 16],
        "пункты МТС, synth (13)":   [3 / 13, 6 / 13, 12 / 13, 12 / 13, 12 / 13, 11 / 13],
        "synth v1 (12)":            [6 / 12, 10 / 12, 8 / 12, 10 / 12, 7 / 12, 9 / 12],
        "v2hard (8)":               [0 / 8, 8 / 8, 7 / 8, 7 / 8, 7 / 8, 7 / 8],
    }
    labels = {
        "doc3 витрина (16)": ["3", "11", "14–16", "16", "16", "12–14"],
        "пункты МТС, synth (13)": ["3", "6", "12", "12", "12", "11"],
        "synth v1 (12)": ["6", "10", "7–9", "10", "7", "9"],
        "v2hard (8)": ["0", "8", "7", "7", "7", "7"],
    }
    fig, axes = plt.subplots(1, 4, figsize=(13.2, 3.9), dpi=150, sharey=True)
    for ax, t in zip(axes, targets):
        vals = data[t]
        bars = ax.bar(range(len(variants)), vals, color=PAL, width=0.72)
        for b, lab in zip(bars, labels[t]):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.02, lab, ha="center", fontsize=9, color="#1f2328")
        ax.set_title(t, fontsize=10.5, weight="bold"); ax.set_ylim(0, 1.15)
        ax.set_xticks(range(len(variants))); ax.set_xticklabels(variants, rotation=35, ha="right", fontsize=8)
        ax.set_yticks([0, 0.5, 1.0]); ax.set_yticklabels(["0", "50%", "100%"])
        ax.axhline(1.0, color="#e5e7eb", lw=0.8)
    axes[0].set_ylabel("recall по голду")
    fig.suptitle("Recall по размеченным дефектам (EXP-13…15, база it6; GPT — диапазон двух прогонов)", fontsize=11)
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


def chart_noise(path: Path) -> None:
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(13.2, 3.7), dpi=150, gridspec_kw={"width_ratios": [1.15, 1]})
    v = ["Детерм. слой", "Промпт GPT-5.5", "Конвейер GPT-5.5", "+ энтропия", "+ логит-зонд", "120b до NA", "120b с NA"]
    n = [0, 15, 9, 7, 7, 11, 7]
    bars = a1.bar(range(len(v)), n, color=PAL[:6] + ["#059669"], width=0.7)
    for b, val in zip(bars, n):
        a1.text(b.get_x() + b.get_width() / 2, val + 0.3, str(val), ha="center", fontsize=9)
    a1.set_xticks(range(len(v))); a1.set_xticklabels(v, rotation=35, ha="right", fontsize=8)
    a1.set_title("Шум: находок medium+ на заведомо чистом документе", fontsize=10.5, weight="bold"); a1.set_ylim(0, 18)

    th = [0, 2, 4, 6, 8]; rs = [11, 11, 10, 9, 5]; rd = [12, 12, 12, 8, 6]; nz = [18, 14, 14, 9, 2]
    a2.plot(th, [x / 12 for x in rs], "o-", color=PAL[1], label="recall synth (из 12)")
    a2.plot(th, [x / 13 for x in rd], "s-", color=PAL[2], label="recall doc3 (из 13)")
    a2b = a2.twinx(); a2b.plot(th, nz, "^--", color=PAL[5], label="шум@clean"); a2b.set_ylabel("шум", color=PAL[5]); a2b.set_ylim(0, 20)
    a2.axvline(4, color="#9ca3af", ls=":", lw=1); a2.text(4.1, 0.08, "θ = 4 принят", fontsize=8, color="#6b7280")
    a2.set_xlabel("порог критика θ"); a2.set_ylim(0, 1.1); a2.set_title("Свип порога критика (EXP-07, it4)", fontsize=10.5, weight="bold")
    h1, l1 = a2.get_legend_handles_labels(); h2, l2 = a2b.get_legend_handles_labels()
    a2.legend(h1 + h2, l1 + l2, fontsize=8, loc="lower left")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


def chart_fixes(path: Path) -> None:
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(13.2, 3.7), dpi=150)
    cats = ["пункты МТС\n(A-класс v3official, из 7)", "doc3\n(из 16)", "слотов UNKNOWN\n(из 27)"]
    before = [3, 13, 10]; after = [7, 16, 0]
    x = range(len(cats))
    a1.bar([i - 0.19 for i in x], before, width=0.36, color="#9ca3af", label="EXP-13 (до фиксов)")
    a1.bar([i + 0.19 for i in x], after, width=0.36, color="#2563eb", label="EXP-14 (после)")
    for i, (b, a) in enumerate(zip(before, after)):
        a1.text(i - 0.19, b + 0.3, str(b), ha="center", fontsize=9); a1.text(i + 0.19, a + 0.3, str(a), ha="center", fontsize=9)
    a1.set_xticks(list(x)); a1.set_xticklabels(cats, fontsize=9); a1.legend(fontsize=8); a1.set_ylim(0, 19)
    a1.set_title("Три дыры, найденные бенчем, и их закрытие", fontsize=10.5, weight="bold")

    steps = ["p0 наивный", "+ цитаты", "+ таксономия", "+ обезличивание", "конвейер v2g", "+ энтропия v2e"]
    rec = [2, 2, 5, 7, 8, 9]; noise = [17, 17, 14, 10, 9, 7]
    a2.plot(range(len(steps)), [r / 12 for r in rec], "o-", color=PAL[1], label="recall synth (из 12)")
    a2b = a2.twinx(); a2b.plot(range(len(steps)), noise, "^--", color=PAL[5], label="шум@clean"); a2b.set_ylim(0, 20); a2b.set_ylabel("шум", color=PAL[5])
    a2.set_xticks(range(len(steps))); a2.set_xticklabels(steps, rotation=25, ha="right", fontsize=8); a2.set_ylim(0, 1.05)
    a2.set_title("Лестница: от промпта к конвейеру (Qwen-14B, EXP-02/09)", fontsize=10.5, weight="bold")
    h1, l1 = a2.get_legend_handles_labels(); h2, l2 = a2b.get_legend_handles_labels(); a2.legend(h1 + h2, l1 + l2, fontsize=8, loc="center right")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


def chart_models(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(13.2, 3.5), dpi=150)
    models = ["Qwen2.5-14B\n(ночь, it3)", "gpt-oss-120b\n(EXP-15, it6)", "GPT-5.5\n(EXP-14, it6)"]
    doc3 = [7 / 13, 13 / 16, 16 / 16]; mts = [None, 11 / 13, 12 / 13]; labels_doc3 = ["7/13", "12–14/16", "14–16/16"]; labels_mts = ["—", "11/13", "12/13"]
    x = range(3)
    b1 = ax.bar([i - 0.19 for i in x], doc3, width=0.36, color=PAL[1], label="doc3 (целевой документ)")
    b2 = ax.bar([i + 0.19 for i in x], [m or 0 for m in mts], width=0.36, color=PAL[3], label="пункты МТС (synth v3official)")
    for b, l in zip(b1, labels_doc3): ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.02, l, ha="center", fontsize=9)
    for b, l in zip(b2, labels_mts): ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.02, l, ha="center", fontsize=9)
    ax.set_xticks(list(x)); ax.set_xticklabels(models, fontsize=9.5); ax.set_ylim(0, 1.15); ax.set_yticks([0, 0.5, 1]); ax.set_yticklabels(["0", "50%", "100%"])
    ax.legend(fontsize=8, loc="upper left"); ax.set_title("Лестница моделей на одном конвейере: модель — параметр, а не архитектура", fontsize=10.5, weight="bold")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)


# ---------------------------------------------------------------- слайды
def _title(s, text, size=26):
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = INK


def _note(s, text, y=6.85):
    n = s.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.2), Inches(0.5)); n.text_frame.word_wrap = True
    p = n.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(11); p.font.color.rgb = MUTED


def _bullets_box(s, bullets, x, y, w, h, size=16):
    body = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = body.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        level = 1 if b.startswith("  ") else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.text = ("• " if level == 0 else "– ") + b.strip(); p.level = level
        p.font.size = Pt(size if level == 0 else size - 2); p.font.color.rgb = INK if level == 0 else MUTED; p.space_after = Pt(5)
    return body


def add_title(prs, title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.4)); tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = INK
    if subtitle:
        tb2 = s.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.6)); tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = MUTED
    return s


def add_bullets(prs, title, bullets, note="", size=17):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _title(s, title)
    _bullets_box(s, bullets, 0.7, 1.3, 11.9, 5.4, size=size)
    if note: _note(s, note)
    return s


def add_stage(prs, title, what, why, evidence, example, note=""):
    """Слайд этапа: что делает / зачем / чем доказано / пример находки."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _title(s, title)
    cols = [("Что делает", what, 0.6), ("Зачем", why, 4.7), ("Чем доказано", evidence, 8.8)]
    for head, items, x in cols:
        h = s.shapes.add_textbox(Inches(x), Inches(1.25), Inches(3.9), Inches(0.4))
        p = h.text_frame.paragraphs[0]; p.text = head; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT
        _bullets_box(s, items, x, 1.65, 3.9, 3.2, size=13)
    ex = s.shapes.add_textbox(Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.75)); ex.text_frame.word_wrap = True
    p = ex.text_frame.paragraphs[0]; p.text = "Пример находки на документах МТС"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = ACCENT
    for line in example:
        q = ex.text_frame.add_paragraph(); q.text = line; q.font.size = Pt(12.5); q.font.color.rgb = INK; q.space_after = Pt(3)
    if note: _note(s, note)
    return s


def add_image(prs, title, img, note="", top=1.15, width=12.3):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _title(s, title)
    s.shapes.add_picture(str(img), Inches((13.333 - width) / 2), Inches(top), width=Inches(width))
    if note: _note(s, note)
    return s


def add_table(prs, title, header, rows, note="", col_widths=None, font=12):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _title(s, title)
    n_rows, n_cols = len(rows) + 1, len(header)
    shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.38) * n_rows)
    table = shape.table
    for j, h in enumerate(header):
        c = table.cell(0, j); c.text = h; c.text_frame.paragraphs[0].font.size = Pt(font); c.text_frame.paragraphs[0].font.bold = True
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = table.cell(i, j); c.text = str(val); c.text_frame.paragraphs[0].font.size = Pt(font)
    if col_widths:
        for j, w in enumerate(col_widths): table.columns[j].width = Inches(w)
    if note: _note(s, note, y=6.75)
    return s


def main() -> int:
    chart_architecture(CH / "architecture.png"); chart_pipeline(CH / "pipeline.png"); chart_recall(CH / "recall.png")
    chart_noise(CH / "noise.png"); chart_fixes(CH / "fixes.png"); chart_models(CH / "models.png")

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # 1
    add_title(prs, "TZ Review — предварительное ревью ТЗ на потоки и витрины данных",
              "Кейс МТС NET · промежуточная версия 04.09.2026 · github.com/FD-AX/MWS\n"
              "Что уже есть: конвейер с доказанными этапами, self-hosted модель, контур с историей и мониторингом, демо на документах МТС")

    # 2
    add_bullets(prs, "Проблема: что доживает до разработки", [
        "Аналитик NET пишет ТЗ на поток или витрину; разработчик и тестировщик вычитывают вручную, вопросы всплывают после передачи",
        "Три класса проблем — все три есть в тестовых документах МТС:",
        "  пропуск: не задан ключ витрины; нет правила для NULL; нет расписания и SLA; NOT NULL/NULLABLE не указан ни для одного поля",
        "  неоднозначность: «FIELD_IMEI (последний по времени за месяц)» — по какому полю времени и из какой таблицы?",
        "  противоречие: «Способ загрузки: Инкремент» и «Обновление: только полная перезагрузка месяца» в одном документе",
        "Кейсодатель: 4 зоны внимания, 8 обязательных пунктов документации, проверка на новом документе типа «витрина», метрика — полезность замечаний",
    ])

    # 3
    add_bullets(prs, "Постановка и наши требования к себе", [
        "Вход: ТЗ (markdown, txt, docx, pdf). Выход: места, требующие уточнения — цитата · почему · что уточнить · критичность; покрытие чеклиста; вердикт",
        "Решение не заменяет аналитика: светофор и вопросы, а не правки",
        "Требования, которые превратились в архитектуру:",
        "  замечание без привязки к тексту — не замечание → программная верификация цитат",
        "  обезличивание (TABLE_*, FIELD_*, CLUSTER) — не дефект → правило в промптах и в критике",
        "  8 пунктов кейсодателя проверяются всегда → отдельные слоты, критик не вправе их резать",
        "  качество измеряется, а не оценивается → голды, синтетика, бенч, журнал экспериментов",
    ])

    # 4
    add_bullets(prs, "Как мы доказываем: данные и метрики", [
        "Голд: 3 обезличенных документа МТС, 36 размеченных дефектов с кодами таксономии (A пропуск, B неоднозначность, C ссылка, D противоречие, E измеримость, F шаблон) и сложностью; 9 позиций — официальные пункты",
        "Фабрика синтетики: чистая база в официальном шаблоне МТС (20 разделов) + рецепты инъекций с автоматическим голдом; 33 инъекции (v1 обычные, v2hard экспертные, v3official — по 8 пунктам); 7 итераций вычитки",
        "Метрики: recall по классам и сложности, шум на заведомо чистом документе, anchoring (доля верифицированных цитат), разброс между прогонами, цена (вызовы, токены)",
        "Бенч «варианты × цели» → матрица ошибок «вариант × класс»: слепые зоны (не ловит никто) и уникальные вклады (ловит один)",
        "Правило журнала: файл эксперимента с вопросом и ожиданием — до прогона; 16 экспериментов за трое суток",
    ], note="Правило проекта: любое решение = «было → стало» цифрами + команда воспроизведения. experiments/README.md")

    # 5
    add_table(prs, "Гипотезы и что с ними стало", ["Гипотеза", "Эксперимент", "Результат", "Вердикт"], [
        ["H1 Граф сущностей из заглушек обезличивания", "EXP-01", "synth 2/12 → 5/12, doc3 +critical, 0 ложных, 0 вызовов LLM", "в прод"],
        ["Лестница промпта: цитаты → таксономия → обезличивание", "EXP-02", "таксономия +3 recall; обезличивание +2 и −40% шума; цитаты recall не двигают, но питают верификацию", "выводы сняты"],
        ["Длина контекста", "EXP-03", "обрыв recall ×3 на ~28k символов; A-класс умирает первым; anchoring 100% → 47%", "чанкование — план"],
        ["H5 «Компиляция ТЗ» (план реализации с журналом допущений)", "EXP-05", "+2 на doc3 в ансамбле; самый чистый проход (шум 5)", "в прод (v3)"],
        ["Порог критика", "EXP-07", "θ=4: recall держится, хвост −40%; шум порогом не лечится", "принят; промпт критика — H11"],
        ["Подход D: семантическая энтропия", "EXP-09/14", "doc3 7/13 → 11/13 (it3); 16/16 (it6) при шуме 7", "в прод"],
        ["H10 Логит-зонд вместо сэмплирования", "EXP-11/13/14", "×0.1 цены; B-класс 9/9; единственный ловил анафору; doc3 16/16", "в прод"],
        ["Официальные пункты МТС как слоты + защита от критика", "EXP-13/14", "A-класс v3official 3/7 → 7/7", "в прод"],
        ["Исход «не применимо» у слота", "EXP-16", "шум 120b 11 → 7 без потери пунктов МТС", "в прод"],
        ["Модель — параметр (лестница Qwen-14B / gpt-oss-120b / GPT-5.5)", "ночь, EXP-15", "модель решает сильнее архитектуры, но детерминированные слои модели не видят", "лестница"],
    ], col_widths=[3.9, 1.4, 5.3, 1.5], font=11)

    # 6
    add_image(prs, "Почему не «один промпт»: лестница и три дыры, найденные бенчем", CH / "fixes.png",
              note="Промпт-бейзлайн GPT-5.5: слеп к официальному шаблону (0 из 6) и промахивает лёгкое (3 из 9). Три дыры EXP-13: критик резал верные находки по пунктам МТС (score 0), треть слотов чеклиста терялась молча, табличные цитаты не верифицировались.", top=1.25, width=12.6)

    # 7
    add_image(prs, "Конвейер: восемь этапов и что каждый доказал", CH / "pipeline.png", top=1.6, width=12.8,
              note="Серые этапы — без LLM. Дальше — по слайду на этап: что делает, зачем, чем доказано, пример находки.")

    # 8..15 — этапы
    add_stage(prs, "Этап 1. Детерминированный слой",
              ["Проверка официального шаблона МТС: 20 разделов на месте, пустой раздел допустим только с «не применимо»",
               "Языковые паттерны: TBD, «и т. д.», «при необходимости», «оперативно»",
               "Обязательные разделы (источник, структура, алгоритм, регламент, контроль качества)"],
              ["Бесплатно, precision ≈ 1: работает до модели и без модели",
               "Даёт гарантированный пол качества на любой модели заказчика",
               "Официальное правило кейсодателя (п.5) проверяется буквально"],
              ["F-класс (шаблон) 6/6 при шуме 0 (EXP-13, EXP-14)",
               "Промпт-бейзлайн на тех же дефектах — 0/6",
               "Не отдаётся критику: его находки нельзя срезать"],
              ["template:official_missing (doc3): отсутствуют разделы Решаемая проблема, Data Catalog, DDL, Пример данных, FAQ, История изменений — по правилам кейсодателя разделы сохраняются или помечаются «не применимо»."])

    add_stage(prs, "Этап 2. Граф сущностей (H1)",
              ["Обезличивание = бесплатная разметка: TABLE_*, FIELD_*, DAG_* — сущности",
               "Поля-фантомы: используются в алгоритме, но не описаны ни в одной структуре",
               "Дрейф имён (TABLE_AGG_LOAD_LOG ↔ TABLE_AGG_LOG_LOAD), битые ссылки «см. раздел N»"],
              ["Класс C (ссылки в пустоту) LLM ловит нестабильно, граф — всегда",
               "0 вызовов модели, не зависит от длины документа",
               "Даёт критичные находки там, где человек не сверит 5 таблиц"],
              ["EXP-01: synth 2/12 → 5/12, doc3 +1 critical, 0 ложных",
               "EXP-13: единственный, кто держит C-класс на всех моделях",
               "EXP-14: критик резал граф-находку → граф выведен из-под критика"],
              ["graph:undefined_field (doc3, critical): поля FIELD_CELL_GEN, FIELD_CELL_ID, FIELD_IMEI, FIELD_IMSI, FIELD_LAC, FIELD_TIME_STAMP используются в алгоритме, но не описаны ни в одной структуре данных документа."])

    add_stage(prs, "Этап 3. Чеклист полноты: 27 слотов",
              ["Бинарные вопросы по домену DWH: инкремент, late-arriving, NULL, дедупликация, ключи, SLA, сбои, объёмы",
               "8 «Основных моментов» МТС — отдельные слоты (SER, CAT, NUL-03, FIL-02, LOC-02/03, REF, KEY-02)",
               "Исходы OK (с цитатой) / MISSING / UNCLEAR / NA («не применимо»); батчи по 5, добор неотвеченных"],
              ["«Оцени качество» не работает; «есть ли X, покажи цитату» — работает (PR-Agent)",
               "Требования кейсодателя проверяются каждый раз, не по настроению модели",
               "NA убирает ложные пробелы: нет Kafka — вопрос о кластере Kafka не пробел"],
              ["EXP-13 → EXP-14: пункты МТС 3/7 → 7/7 после защиты от критика",
               "Схема ответа теряла OK-статусы (10/27 UNKNOWN) → 0 после фикса; зонды получили свои слоты",
               "EXP-16: NA — шум 120b 11 → 7 без потери пунктов"],
              ["checklist:NUL-03 (doc3, high): в структуре витрины перечислены поля, типы и описания, но ни для одного поля не указан признак NOT NULL / NULLABLE (пункт МТС №3). Что уточнить: для каждого поля витрины — обязательность."])

    add_stage(prs, "Этап 4. Согласованность разделов",
              ["Межсекционная проверка: поле без источника, дрейф термина, противоречие разделов",
               "Смотрит на документ целиком, а не на слот",
               "Категории: doc:contradiction, doc:terminology, doc:completeness"],
              ["Противоречия — то, что хуже всего видит человек, читающий по разделам",
               "Чеклист слот за слотом их не найдёт: каждый раздел по отдельности «в порядке»",
               "Класс D в голде МТС — 15 дефектов из 69"],
              ["EXP-14: D-класс 10/11 у v2e, 14/15 в EXP-13 у v2g",
               "Ловит G3-1 (инкремент vs полная перезагрузка), G1-1 (задержка <1 мин vs ≈0 сек)",
               "Уникальный вклад в EXP-10: G2-5 (часы vs минуты в сдвигах поясов)"],
              ["doc:contradiction (doc3, critical): «Способ загрузки: Инкремент», а в требованиях к агрегату «Обновление: только полная перезагрузка месяца (без upsert)». Что уточнить: инкремент или полная перезапись партиции?"])

    add_stage(prs, "Этап 5. Взгляд разработчика",
              ["Персона: «ты должен начать реализацию сегодня — где придётся угадывать?»",
               "Вопросы, блокирующие код, с ценой ошибки (why); максимум 8, без общих советов",
               "Запрещены вопросы о настоящих именах заглушек"],
              ["Это ровно та ситуация, которую кейс называет болью: разработчик угадывает",
               "Даёт дефекты вне рубрики — краевые случаи логики (A3, B9), которых нет в чеклисте",
               "Формулирует вопрос так, как его задали бы аналитику"],
              ["Уникальные находки класса A3/B9 на doc3: TAC не найден → vendor?, «последний по времени» — по какому полю",
               "EXP-02: без таксономии и правила обезличивания — generic-каша (шум 17)",
               "EXP-16: правило заглушек убрало «как называется столбец FIELD_IMSI»"],
              ["dev_question (doc3, medium, score 5): «substring(imei,1,8) = tac → vendor_name». Что делать, если полученный TAC отсутствует в TABLE_DEVICE_REF: NULL, «UNKNOWN» или иной fallback? Неправильный выбор исказит агрегат и контроль качества."])

    add_stage(prs, "Этап 6. Неоднозначность как число: энтропия и логит-зонд",
              ["Слоты со статусом OK перечитываются: 5 сэмплов с температурой → кластеры ответов → энтропия (v2e)",
               "Или один вызов на слот: P(YES)/P(NO) из top-logprobs, бинарная энтропия (v2l, H10)",
               "Расходятся ответы = документ читается по-разному"],
              ["Чеклист говорит «есть», а два разработчика прочитают по-разному — это класс B",
               "Измерение вместо мнения: порог калибруется по голду, воспроизводимо",
               "Логиты в 10 раз дешевле сэмплирования и детерминированы"],
              ["EXP-09: doc3 7/13 → 11/13 (+4) при шуме 10 → 7",
               "EXP-13/14: логит-зонд единственный ловил анафору «Он определяется…»; B-класс 9/9",
               "EXP-14: v2e и v2l — doc3 16/16, anchoring 96–100%"],
              ["uncertainty (doc3): слот SLA-01 «расписание задано?» — чеклист ответил OK по «Ежемесячно (1 раз в месяц)», сэмплы разошлись (время запуска? какой DAG?) → находка G3-9: нет времени запуска и SLA, оркестрация не описана."],
              note="Требует logprobs у модели: GPT-4.1 / vLLM — да, ollama — нет → план: vLLM для gpt-oss-120b.")

    add_stage(prs, "Этап 7. Верификация цитат",
              ["Цитата обязана быть подстрокой документа (после нормализации); не нашлась — находка отбрасывается",
               "Нашлась в другом разделе — переанкоривается",
               "Вторая ступень для табличных цитат без «|» и с другими переносами"],
              ["Галлюцинированная цитата — главный способ потерять доверие аналитика",
               "Даёт метрику anchoring: доля находок, которым можно верить дословно",
               "Единственная защита, когда модель деградирует на длинных документах"],
              ["EXP-03: на 28k символов anchoring 100% → 47% — verify срезал галлюцинации",
               "EXP-13: 5 из 21 находок отброшены из-за табличных цитат → вторая ступень",
               "EXP-14: anchoring 96–100% у прод-вариантов на doc3"],
              ["dropped (doc3 v2, 120b): «FIELD_USERS_CNT является результатом агрегирования… включать его в группировку противоречит SQL» — цитата не найдена в документе → находка не показана аналитику (в отчёте видна как «отброшено верификацией»)."])

    add_stage(prs, "Этап 8. Критик",
              ["Судья видит все находки разом (Greptile), ставит score 0–10 сравнительно, группирует дубли",
               "Порог θ=4 по свипу; официальные пункты МТС и детерминированные проходы порогу не подлежат",
               "Score 0 за вопросы о настоящих именах заглушек и за «информация уже есть в документе»"],
              ["Генерация даёт кандидатов, фильтрация — ответственность отдельного шага",
               "Аналитику нужны 5 точных находок, а не 15 «на всякий случай»",
               "Но фильтр не вправе спорить с требованиями кейсодателя"],
              ["EXP-07: θ=4 — хвост находок −40% (60 → 36), recall doc3 12/13 держится",
               "EXP-13: критик ставил 0.0 верным MISSING по Data Catalog, Kafka, фильтрам → защита слотов",
               "EXP-14: пункты МТС 3/7 → 7/7; шум 9 → 7 у v2e/v2l"],
              ["rejected (doc3, 120b): checklist:INC-01 score 0 — «нет механизма инкремента», при том что документ явно описывает полную перезапись партиции; критик убрал ложный пробел. Из 25 кандидатов в отчёт прошло 16."])

    # 16
    add_image(prs, "Результаты: recall по размеченным дефектам", CH / "recall.png", top=1.2, width=12.6,
              note="Промпт слеп к шаблону и промахивает лёгкое; конвейер закрывает все 8 пунктов МТС; зонды неоднозначности дают 16/16 на целевом документе. Разброс GPT между прогонами ±2 → числа как диапазон, консенсус двух прогонов.")

    # 17
    add_image(prs, "Шум и порог критика", CH / "noise.png", top=1.2, width=12.6,
              note="Шум порогом не лечится (EXP-07): при θ=6 recall doc3 падает 12 → 8. Лечится правилами: обезличивание (−40%), «не применимо» (11 → 7 на 120b), запрет вопросов о заглушках. Трижды «шум» на чистой синтетике оказывался реальными дефектами базы → правило: читать сырые выводы глазами.")

    # 18
    add_image(prs, "Лестница моделей: разрыв self-hosted 120b и GPT-5.5", CH / "models.png", top=1.2, width=12.2,
              note="Разрыв на doc3 — 4 «выводных» дефекта (по какому полю времени последний IMEI, TAC не найден, бэкфилл, границы UTC); у GPT их добирают зонды. Путь: vLLM + logprobs на 120b. На пунктах МТС разрыв 1 дефект (11 vs 12 из 13). H100 $3.49/ч ≈ $0.4/документ.")

    # 19
    add_image(prs, "Что уже работает: контур и демо", CH / "architecture.png", top=1.2, width=12.2,
              note="Проверено на документах МТС: PDF → секции (doc1 11, doc2 8, doc3 14) → очередь → воркер на gpt-oss-120b → отчёт → история; doc3 14/16, doc1 12/12 по голду. UI: прогресс по этапам с оценкой времени, история версий документа, 👍/👎 по находке. RabbitMQ (DLQ), Postgres (документы, ревью, находки, фидбек), Prometheus/Grafana.")

    # 20
    add_table(prs, "Риски AI-решения и как закрываем", ["Риск", "Проявление у нас", "Мера / статус"], [
        ["Галлюцинация цитаты", "модель «цитирует» отсутствующее", "верификация цитат; anchoring 96–100%"],
        ["Фильтр режет правду", "критик ставил 0 верным пунктам МТС (EXP-13)", "защищённые слоты и детерминированные проходы (EXP-14: 3/7 → 7/7)"],
        ["Молчаливая потеря ответов", "10/27 слотов UNKNOWN (EXP-13)", "схема, добор, логирование; UNKNOWN на дашборде — 0"],
        ["Разброс модели", "GPT ±2, 120b ±1–2 дефекта", "консенсус двух прогонов; числа как диапазон"],
        ["Ложные срабатывания на обезличивании", "«как называется столбец FIELD_IMSI» (120b)", "правило в промптах + score 0 в критике (EXP-16)"],
        ["Синтетика ≠ реальность", "«шум» трижды был правдой", "вычитка сырых выводов, it1–it7; тройная разметка голдов — до финала"],
        ["Слабая self-hosted модель", "пропускает выводные дефекты", "зонды через vLLM/logprobs; reasoning effort — измерить"],
        ["Длинные документы", "обрыв на ~28k символов (EXP-03)", "чанкование по секциям — план"],
    ], col_widths=[3.0, 4.2, 4.9], font=11)

    # 21
    add_bullets(prs, "План до финала (07.09)", [
        "vLLM вместо ollama для gpt-oss-120b → логит-зонд и энтропия на self-hosted модели; цель doc3 ≥ 15/16",
        "Тройная разметка голдов документов МТС (второй и третий размётчик), пересчёт всех таблиц",
        "Детерминированные детекторы для оставшихся слепых зон: ссылка на раздел при ненумерованных заголовках, объявленная таблица без структуры, тип против значений; INC-01 → «не применимо» при полной перезаписи",
        "Дашборд качества: recall по итерациям рядом с эксплуатационными метриками (bench_runs в Postgres)",
        "Фронтенд аналитика на нашем API (контракт готов, CORS включён); отчёт-обоснование архитектуры по журналу; финальное демо на новом документе кейсодателя",
    ])

    # 22
    add_table(prs, "Команда и распределение задач", ["Участник", "Роль", "Зона ответственности", "Степень участия"], [
        ["Коротков Леонид (@fDAFX), магистрант ИТМО, программа «Искусственный интеллект» (AI Talent Hub)", "AI Engineer — исследование и направление",
         "генерация и отбор идей: гипотезы, архитектурные решения, план экспериментов, выбор моделей и метрик; методология оценки", "ведущая, постоянно"],
        ["Рыженко Виктор (@RyzhenkoViktor)", "AI Engineer — реализация",
         "код и инфраструктура: конвейер ревью и тесты, бенч и фабрика синтетики, compose-контур, интеграция self-hosted модели", "частичная, по задачам реализации"],
        ["Сокол Тимофей (@xcmgg)", "AI Product — процессы и продукт",
         "процесс команды и курс развития, продуктовая проработка (гипотеза, MVP, критерии, риски), фронтенд аналитика, участие в коде", "высокая, постоянно"],
    ], note="Степень участия — на 04.09; в финальной версии уточняется по фактическому вкладу.", col_widths=[2.9, 2.5, 4.9, 1.8])

    prs.save(OUT)
    print("saved", OUT, "slides:", len(prs.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
